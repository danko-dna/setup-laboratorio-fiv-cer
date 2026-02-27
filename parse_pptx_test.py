import pptx
import sys

def hex_color(color):
    if color is None or not hasattr(color, 'rgb') or color.type != 1:
        return 'None'
    return str(color.rgb)

def analyze_pptx(file_path):
    print(f"Analyzing {file_path}...")
    try:
        prs = pptx.Presentation(file_path)
    except Exception as e:
        print(f"Error: {e}")
        return

    for s_idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {s_idx+1} ---")
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"Table: {len(table.rows)}x{len(table.columns)}")
                for r_idx, row in enumerate(table.rows):
                    row_data = []
                    for cell in row.cells:
                        text = cell.text.replace('\n', ' ')
                        
                        # Try to get cell background color
                        fill = cell.fill
                        color_hex = "None"
                        if fill.type == pptx.enum.dml.MSO_FILL.SOLID:
                            try:
                                color_hex = fill.fore_color.rgb
                            except:
                                pass
                        
                        # Try to get text color from the first run
                        text_color_hex = "None"
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                try:
                                    if r.font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.RGB:
                                        text_color_hex = r.font.color.rgb
                                        break
                                except:
                                    pass
                            if text_color_hex != "None":
                                break

                        row_data.append(f"[{text} (b:{color_hex}, f:{text_color_hex})]")
                    print(f"Row {r_idx}: " + " | ".join(row_data))
            elif shape.has_text_frame:
                print(f"Text: {shape.text.strip().replace(chr(11), ' ')}")

if __name__ == "__main__":
    analyze_pptx(sys.argv[1])
