import pptx
import sys

def parse_pptx_simple(file_path):
    print(f"Reading {file_path}:\n" + "="*40)
    prs = pptx.Presentation(file_path)
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}")
        for shape in slide.shapes:
            print(f"- Shape: {shape.name} (Type: {shape.shape_type})")
            if shape.has_table:
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    row_data = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                    print(f"Row {r_idx}: {row_data}")
                print("-" * 20)
            elif shape.has_text_frame:
                print(f"Text: {shape.text.replace(chr(11), ' ').strip()}")

if __name__ == "__main__":
    for arg in sys.argv[1:]:
        parse_pptx_simple(arg)
